from playwright.sync_api import sync_playwright
import os
import requests
from pathlib import Path
from typing import Optional
from io import BytesIO
import docx
import PyPDF2
import openpyxl
import pptx

class PlaywrightCrawler:
    def __init__(self, headless: bool = True):
        self.headless = headless

    def fetch_page(self, url: str) -> str:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=self.headless)
            page = browser.new_page()
            page.goto(url)
            content = page.content()
            browser.close()
            return content
        
    def get_download_links(self, url: str) -> list:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=self.headless)
            page = browser.new_page()
            page.goto(url)
            # Find all <a> tags with href that looks like a file link
            links = page.eval_on_selector_all(
                "a[href]",
                """elements => elements
                    .map(el => el.href)
                    .filter(href => /\\.(pdf|docx?|xlsx?|pptx?)$/i.test(href))
                """
            )
            browser.close()
            return links

    def download_file(self, url: str, dest_folder: str = "./downloads") -> Optional[str]:
        os.makedirs(dest_folder, exist_ok=True)
        local_filename = os.path.join(dest_folder, url.split("/")[-1])
        try:
            with requests.get(url, stream=True) as r:
                r.raise_for_status()
                with open(local_filename, 'wb') as f:
                    for chunk in r.iter_content(chunk_size=8192):
                        f.write(chunk)
            return local_filename
        except Exception as e:
            print(f"Failed to download {url}: {e}")
            return None

    def parse_document(self, file_path: str) -> Optional[str]:
        ext = Path(file_path).suffix.lower()
        try:
            if ext == ".pdf" and PyPDF2:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() or ""
                    return text
            elif ext in [".doc", ".docx"] and docx:
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            elif ext in [".xls", ".xlsx"] and openpyxl:
                wb = openpyxl.load_workbook(file_path, read_only=True)
                text = ""
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                return text
            elif ext in [".ppt", ".pptx"] and pptx:
                prs = pptx.Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                print(f"Unsupported file type or missing parser for: {file_path}")
                return None
        except Exception as e:
            print(f"Failed to parse {file_path}: {e}")
            return None
if __name__ == "__main__":
    crawler = PlaywrightCrawler()
    url = "https://example.com"
    html = crawler.fetch_page(url)
    print(html)